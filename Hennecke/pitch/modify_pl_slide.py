#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Replace the P&L cost-structure slide block with a 2023-vs-2024 YoY version
using Hennecke-OMS real figures (% of value of production)."""
f = "/home/user/Storage/Hennecke/pitch/build_pptx_editable.py"
src = open(f).read()

# ensure MSO_CONNECTOR import
if "MSO_CONNECTOR" not in src:
    src = src.replace("from pptx.enum.shapes import MSO_SHAPE",
                      "from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR")

NEW = r'''# =====================================================================
# SLIDE 2a2 — P&L COST STRUCTURE 2023 vs 2024 (YoY)
# =====================================================================
s = slide()
header(s, "P&L recap — cost structure (2023 vs 2024)",
       "Value of production split into EBITDA and cost categories · year-on-year", "3 / 7")
cats = ["EBITDA", "Materials", "Personnel", "Services", "Lease & rentals", "Other op. costs"]
colA = {"EBITDA":"84BD00","Materials":"F2A100","Personnel":"D9531E","Services":"3C7DA6","Lease & rentals":"8C5A2B","Other op. costs":"BFBFBF"}
txtA = {"EBITDA":INK,"Materials":INK,"Personnel":WHITE,"Services":WHITE,"Lease & rentals":WHITE,"Other op. costs":INK}
d23 = {"EBITDA":7.0,"Materials":46.0,"Personnel":23.1,"Services":19.9,"Lease & rentals":3.8,"Other op. costs":0.3}
d24 = {"EBITDA":3.8,"Materials":38.4,"Personnel":30.3,"Services":22.8,"Lease & rentals":4.2,"Other op. costs":0.5}
deltas = {"EBITDA":"-3.2","Materials":"-7.6","Personnel":"+7.2","Services":"+2.9","Lease & rentals":"+0.4","Other op. costs":"+0.2"}
TOPB, SCALE, BW = 165, 6.0, 130
def draw_bar(x, data):
    y = TOPB; bounds = [y]
    for c in cats:
        h = data[c]*SCALE
        rect(s, x, y, BW, h, fill=colA[c])
        if h >= 17:
            para(tb(s, x, y, BW, h, "m"), f"{data[c]:.1f}%", size=(13 if h >= 28 else 10), color=txtA[c], bold=True, align="c", first=True)
        y += h; bounds.append(y)
    return bounds
BX1, BX2 = 480, 720
b1 = draw_bar(BX1, d23)
b2 = draw_bar(BX2, d24)
for i in range(len(b1)):
    cn = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, PX(BX1+BW), PX(b1[i]), PX(BX2), PX(b2[i]))
    cn.line.color.rgb = RGBColor.from_string("D9DEE6"); cn.line.width = Pt(0.75); cn.shadow.inherit = False
para(tb(s, BX1-20, b1[-1]+8, BW+40, 24), "2023", size=15, color=BLUE, bold=True, align="c", first=True)
para(tb(s, BX2-20, b2[-1]+8, BW+40, 24), "2024", size=15, color=BLUE, bold=True, align="c", first=True)
para(tb(s, 350, TOPB-2, 120, 22), "% of value", size=11, color=MUTED, align="r", first=True)
para(tb(s, 350, TOPB+16, 120, 22), "of production", size=11, color=MUTED, align="r", first=True)
# recap table (manual, fully styled)
rect(s, 960, 190, 590, 300, fill=LIGHT, line=BORDER)
para(tb(s, 982, 200, 260, 26), "Category", size=12, color=GRAY, bold=True, first=True)
para(tb(s, 1250, 200, 96, 26), "2023", size=12, color=GRAY, bold=True, align="c", first=True)
para(tb(s, 1352, 200, 96, 26), "2024", size=12, color=GRAY, bold=True, align="c", first=True)
para(tb(s, 1454, 200, 90, 26), "Δ pp", size=12, color=GRAY, bold=True, align="c", first=True)
rect(s, 982, 230, 562, 1, fill=BORDER)
ry = 244
for c in cats:
    para(tb(s, 982, ry, 268, 28, "m"), [("■ ", False, colA[c], 12), (c, True, INK, 12)], first=True)
    para(tb(s, 1250, ry, 96, 28, "m"), f"{d23[c]:.1f}%", size=12, color=TXT, align="c", first=True)
    para(tb(s, 1352, ry, 96, 28, "m"), f"{d24[c]:.1f}%", size=12, color=TXT, align="c", first=True)
    dv = deltas[c]; dcol = "2E7D4F" if (dv.startswith("-") and c != "EBITDA") else "C8102E"
    para(tb(s, 1454, ry, 90, 28, "m"), dv, size=12, color=dcol, bold=True, align="c", first=True)
    ry += 36
# callout
rect(s, 960, 505, 590, 150, fill=BLUEBG, line=BLUELN)
tfc = tb(s, 982, 517, 546, 130)
para(tfc, [("Fixed-cost absorption is the story: ", True, BLUE, 14), ("personnel rises ", False, TXT, 14), ("+7pp", True, TXT, 14), (" of output (stable in €, heavier on lower volume) while materials fall ", False, TXT, 14), ("-8pp", True, TXT, 14), (" (variable) — net, EBITDA compresses from 7.0% to 3.8%.", False, TXT, 14)], first=True, after=6)
para(tfc, "Read with the project cycle: 2024 value of production is lower as the 2023 WIP build-up unwinds.", size=12, color=GRAY, italic=True)
para(tb(s, 50, 860, 1330, 30), "% of value of production (€52.0m 2023 / €40.7m 2024). 2023 base inflated by WIP build-up (+€19.3m); on a revenue basis EBITDA margin was 7.2% (2023) / 3.5% (2024). Source: Coface reclassified P&L.", size=9.5, color=MUTED, first=True)

'''

start = src.index("# SLIDE 2a2 ")
banner = src.rindex("# ====", 0, start)
end = src.index("# =====================================================================\n# SLIDE 2b — CUSTOMER MAPPING")
src = src[:banner] + NEW + src[end:]
open(f, "w").write(src)
print("P&L slide replaced with YoY version")
