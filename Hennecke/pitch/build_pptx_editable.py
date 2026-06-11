#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Build a fully EDITABLE 5-slide Board pitch (native PowerPoint shapes/text/chart),
not flattened images."""
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION

OUT = "/home/user/Storage/Hennecke/pitch/HenneckeOMS_Board_Pitch_EN.pptx"

# ---- palette ----
BLUE="1F3A5F"; RED="C8102E"; GOLD="F2C200"; GREEN="2E7D4F"
LIGHT="F7F8FA"; BORDER="DFE4EA"; GRAY="666666"; MUTED="999999"
DARK="1A3050"; PANEL="284869"; WHITE="FFFFFF"; LBLUE="9FB2C8"
INK="222222"; TXT="333333"; PINKBG="FDEEF0"; PINKLN="F3C5CC"
GREENBG="EEF6F0"; GREENLN="CFE6D8"; BLUEBG="EEF2F7"; BLUELN="CDD8E6"

EMU = 7620  # EMU per "px" on a 1600x900 design over 13.333x7.5in
def PX(v): return Emu(int(round(v*EMU)))

prs = Presentation()
prs.slide_width = Emu(int(13.333*914400)); prs.slide_height = Emu(int(7.5*914400))
BLANK = prs.slide_layouts[6]

ALIGN = {"l":PP_ALIGN.LEFT, "c":PP_ALIGN.CENTER, "r":PP_ALIGN.RIGHT}
ANCH = {"t":MSO_ANCHOR.TOP, "m":MSO_ANCHOR.MIDDLE, "b":MSO_ANCHOR.BOTTOM}

def slide():
    return prs.slides.add_slide(BLANK)

def rect(s, x, y, w, h, fill=None, line=None, line_w=1.0):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, PX(x), PX(y), PX(w), PX(h))
    sh.shadow.inherit = False
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid(); sh.fill.fore_color.rgb = RGBColor.from_string(fill)
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = RGBColor.from_string(line); sh.line.width = Pt(line_w)
    return sh

def tb(s, x, y, w, h, anchor="t"):
    box = s.shapes.add_textbox(PX(x), PX(y), PX(w), PX(h))
    tf = box.text_frame; tf.word_wrap = True
    tf.vertical_anchor = ANCH[anchor]
    tf.margin_left = Pt(2); tf.margin_right = Pt(2)
    tf.margin_top = Pt(1); tf.margin_bottom = Pt(1)
    return tf

def para(tf, runs, size=14, color=INK, bold=False, align="l", first=False,
         after=4, before=0, italic=False, spacing=1.0):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = ALIGN[align]; p.space_after = Pt(after); p.space_before = Pt(before)
    try: p.line_spacing = spacing
    except Exception: pass
    if isinstance(runs, str):
        runs = [(runs, bold, color, size)]
    for rn in runs:
        text = rn[0]
        b = rn[1] if len(rn) > 1 and rn[1] is not None else bold
        c = rn[2] if len(rn) > 2 and rn[2] is not None else color
        sz = rn[3] if len(rn) > 3 and rn[3] is not None else size
        r = p.add_run(); r.text = text
        f = r.font; f.name = "Arial"; f.size = Pt(sz); f.bold = b; f.italic = italic
        f.color.rgb = RGBColor.from_string(c)
    return p

def header(s, title, subtitle, page=None):
    rect(s, 0, 0, 1600, 8, fill=RED)
    para(tb(s, 50, 30, 1100, 50), title, size=26, color=BLUE, bold=True, first=True)
    para(tb(s, 50, 78, 1200, 30), subtitle, size=13, color=GRAY, first=True)
    para(tb(s, 1150, 26, 400, 22, "t"), "CONFIDENTIAL · For Board discussion",
         size=9.5, color=MUTED, align="r", first=True)
    if page:
        rect(s, 50, 858, 1500, 1, fill=BORDER)
        para(tb(s, 1400, 866, 150, 20), page, size=9.5, color=BLUE, bold=True, align="r", first=True)

def card_title(s, x, y, w, text, accent=BLUE):
    """colored header band + title (editable)."""
    rect(s, x, y, w, 46, fill=accent)
    para(tb(s, x+22, y, w-30, 46, "m"), text, size=15, color=WHITE, bold=True, first=True)

def bullets(tf, items, accent=RED, size=14, color=INK, first_done=False, after=8):
    for i, it in enumerate(items):
        runs = [("•  ", False, accent, size)]
        if isinstance(it, list):
            runs += it
        else:
            runs += [(it, False, color, size)]
        para(tf, runs, size=size, color=color, first=(i == 0 and not first_done), after=after)

# =====================================================================
# SLIDE 0 — COVER
# =====================================================================
s = slide()
rect(s, 0, 0, 1600, 900, fill=BLUE)
rect(s, 0, 0, 1600, 10, fill=RED)
rect(s, 0, 640, 1600, 260, fill=DARK)
para(tb(s, 120, 230, 900, 30), "PRIVATE EQUITY · PROJECT RECAP", size=15, color=LBLUE, bold=False, first=True)
rect(s, 120, 272, 90, 5, fill=RED)
para(tb(s, 116, 300, 1300, 110), "HENNECKE-OMS S.p.A.", size=58, color=WHITE, bold=True, first=True)
para(tb(s, 120, 415, 1200, 50), "Investment opportunity assessment", size=27, color=GOLD, first=True)
rect(s, 122, 500, 778, 2, fill="3A5578")
tfd = tb(s, 120, 525, 1200, 90)
para(tfd, "Italian polyurethane-machinery hub of the Hennecke Group", size=19, color="CDD8E6", first=True, after=4)
para(tfd, "(part of Brückner Group SE) — carve-out / platform evaluation", size=19, color="CDD8E6")
tfb = tb(s, 120, 695, 700, 80)
para(tfb, "Presentation to the Board of Directors", size=18, color=WHITE, bold=True, first=True, after=4)
para(tfb, "June 2026", size=16, color=LBLUE)
tfr = tb(s, 1000, 695, 480, 80, "t")
para(tfr, "Based on FY2023–2024 filed financials", size=14, color=LBLUE, align="r", first=True, after=3)
para(tfr, "Sources: Coface · Chamber of Commerce · public", size=14, color=LBLUE, align="r")
rect(s, 120, 820, 1360, 1, fill="2D4769")
para(tb(s, 120, 832, 1360, 24), "CONFIDENTIAL — For Board discussion only. Not an investment recommendation; subject to specialist Due Diligence.",
     size=11, color="7E93AC", first=True)

# =====================================================================
# SLIDE 1 — OPPORTUNITY
# =====================================================================
s = slide()
header(s, "HENNECKE-OMS S.p.A. — Investment opportunity",
       "Italian PU-machinery hub of the Hennecke Group · potential carve-out / platform deal", "1 / 7")
# left card
rect(s, 50, 125, 720, 700, fill=LIGHT, line=BORDER)
card_title(s, 50, 125, 720, "Company at a glance", BLUE)
tf = tb(s, 74, 188, 672, 360)
bullets(tf, [
    [("Italian production hub of the ", False), ("Hennecke Group", True),
     (" — a world leader in polyurethane (PU) processing machinery.", False)],
    [("HQ & plant in ", False), ("Verano Brianza (MB), Italy", True),
     ("; roots in the historic OMS engineering district (since 1986).", False)],
    [("100% owned by Hennecke GmbH (DE)", True),
     (" — captive subsidiary, subject to direction & coordination.", False)],
    [("Turnkey lines for flexible & rigid PU foam: slabstock, sandwich panels, refrigeration, automotive, footwear.", False)],
    [("~166 employees (2024); ", False), ("ISO 9001 & 14001", True), (" certified.", False)],
    [("Global installed base — recurring aftermarket potential.", False)],
], accent=RED, size=15, after=9)
# group context box
rect(s, 74, 585, 672, 200, fill="EEF2F7", line=BLUELN)
tfg = tb(s, 96, 600, 628, 175)
para(tfg, "Group context", size=15, color=BLUE, bold=True, first=True, after=6)
para(tfg, [("The parent Hennecke Group was acquired by ", False, TXT),
           ("Brückner Group SE", True, TXT),
           (" (closing 5 Jan 2026), from a fund advised by Capvis.", False, TXT)], size=14, after=4)
para(tfg, "Hennecke = ~680 employees, ~15 companies, two brands (Hennecke Polyurethane Technology & Hennecke-OMS).", size=14, color=TXT, after=4)
para(tfg, "Upstream ownership inferred from public sources.", size=12, color=MUTED)
# right KPI tiles
para(tb(s, 800, 128, 600, 30), "Key figures 2024", size=19, color=BLUE, bold=True, first=True)
def kpi(x, y, w, h, label, value, sub, subcolor=RED, vsize=34):
    rect(s, x, y, w, h, fill=WHITE, line=BORDER)
    tf = tb(s, x+20, y+12, w-30, h-20, "t")
    para(tf, label, size=14, color=MUTED, first=True, after=4)
    para(tf, value, size=vsize, color=BLUE, bold=True, after=4)
    para(tf, sub, size=13, color=subcolor)
kpi(800, 175, 360, 120, "Revenue 2024", "€ 44.2M", "-12% YoY")
kpi(1190, 175, 360, 120, "EBITDA 2024", "€ 1.5M", "3.5% margin (vs 7.2% in 2023)")
kpi(800, 315, 360, 120, "Liquidity vs bank debt", "€ 10.4M", "cash · vs € 6.9M M/L bank debt", subcolor=GREEN, vsize=30)
kpi(1190, 315, 360, 120, "Shareholders' equity", "€ 6.0M", "thin — 15.5% of total assets", subcolor=MUTED, vsize=30)
# deal angle
rect(s, 800, 470, 750, 355, fill=BLUE)
para(tb(s, 824, 484, 700, 34), "Deal angle", size=19, color=WHITE, bold=True, first=True)
rect(s, 824, 526, 702, 1, fill="3A5578")
tfd = tb(s, 824, 545, 702, 270)
bullets(tfd, [
    [("Carve-out", True, WHITE), (" of the Italian entity, or acquisition of the entire Hennecke Group as a platform.", False, "E9EEF5")],
    [("Light balance sheet: low debt, high liquidity, working capital self-financed by customer advances (€10.9M).", False, "E9EEF5")],
    [("Shares currently ", False, "E9EEF5"), ("pledged", True, WHITE), (" — to be released on sale.", False, "E9EEF5")],
    [("Coface credit score 5/10 (medium-high risk).", False, "E9EEF5")],
], accent=GOLD, size=16, color="E9EEF5", after=12)

# =====================================================================
# SLIDE 2 — FINANCIALS
# =====================================================================
s = slide()
header(s, "Financial profile 2023–2024",
       "Cyclical & project-based business · thin margins · self-financed working capital", "2 / 7")
# left card with native chart
rect(s, 50, 125, 720, 700, fill=LIGHT, line=BORDER)
para(tb(s, 74, 140, 500, 28), "2023 vs 2024 (€ million)", size=17, color=BLUE, bold=True, first=True)
cd = CategoryChartData()
cd.categories = ["Revenue", "Value of prod.", "EBITDA"]
cd.add_series("2023", (50.4, 52.0, 3.6))
cd.add_series("2024", (44.2, 40.7, 1.5))
gframe = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, PX(70), PX(180), PX(680), PX(470), cd)
chart = gframe.chart
chart.has_legend = True; chart.legend.position = XL_LEGEND_POSITION.TOP
chart.legend.include_in_layout = False
chart.plots[0].gap_width = 80
plot = chart.plots[0]
plot.has_data_labels = True
plot.data_labels.number_format = '0.0'; plot.data_labels.number_format_is_linked = False
plot.data_labels.font.size = Pt(11); plot.data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
chart.series[0].format.fill.solid(); chart.series[0].format.fill.fore_color.rgb = RGBColor.from_string(LBLUE)
chart.series[1].format.fill.solid(); chart.series[1].format.fill.fore_color.rgb = RGBColor.from_string(BLUE)
cax = chart.value_axis; cax.has_major_gridlines = True
cax.tick_labels.font.size = Pt(10)
chart.category_axis.tick_labels.font.size = Pt(11)
# margin callout
rect(s, 74, 670, 666, 130, fill=PINKBG, line=PINKLN)
tfm = tb(s, 94, 682, 626, 110)
para(tfm, "EBITDA margin: 7.2% → 3.5%", size=16, color=RED, bold=True, first=True, after=5)
para(tfm, [("Drop driven by the WIP swing: change in work-in-progress went from ", False, TXT),
           ("+€19.3M (2023)", True, TXT), (" to ", False, TXT), ("-€0.5M (2024)", True, TXT),
           (" — a timing effect of long projects, not structural deterioration.", False, TXT)], size=14)
# right top: balance sheet highlights
rect(s, 800, 125, 750, 320, fill=WHITE, line=BORDER)
para(tb(s, 824, 138, 700, 28), "Balance-sheet highlights", size=17, color=BLUE, bold=True, first=True)
rect(s, 824, 172, 702, 1, fill=BORDER)
tfb = tb(s, 824, 185, 702, 250)
rows = [
    [("€ 10.4M", True), (" cash  ·  vs ", False), ("€ 6.9M", True), (" M/L bank debt (new 2024)", False)],
    [("€ 10.9M", True), (" customer advances fund the working capital", False)],
    [("€ 6.0M", True), (" equity — thin (15.5% of total assets)", False)],
    [("€ 12.0M", True), (" intangibles (goodwill) — impairment watch", False)],
    [("€ 29.5M", True), (" total debt — largely operating & intercompany", False)],
    [("Working capital improving: inventory ", False), ("DIO 119 → 66 days", True)],
]
for i, r in enumerate(rows):
    para(tfb, r, size=15, color=INK, first=(i == 0), after=10)
# right bottom: what the numbers say
rect(s, 800, 470, 750, 355, fill=BLUEBG, line=BLUELN)
para(tb(s, 824, 482, 700, 28), "What the numbers say", size=17, color=BLUE, bold=True, first=True)
rect(s, 824, 516, 702, 1, fill=BLUELN)
tfn = tb(s, 824, 528, 702, 290)
bullets(tfn, [
    [("Revenue -12%, value of production -22%: the gap is entirely the WIP timing effect (long-cycle projects).", False)],
    [("Margin compressed by lower fixed-cost absorption (personnel stable in value, higher in incidence).", False)],
    [("Net loss of ", False), ("€ 0.24M", True), (" in 2024 (vs +€0.29M in 2023).", False)],
    [("Normalised EBITDA", True), (" (cycle + intercompany adjusted) is the right basis for valuation, not the 2024 point figure.", False)],
], accent=BLUE, size=15, color=TXT, after=11)
para(tb(s, 50, 866, 900, 20), "Source: Coface Full Report (filed FY2023–2024). Amounts in € million.",
     size=9.5, color=MUTED, first=True)

# =====================================================================
# SLIDE 2a2 — P&L COST STRUCTURE 2023 vs 2024 (YoY)
# =====================================================================
s = slide()
header(s, "P&L recap — cost structure (2023 vs 2024)",
       "Value of production split into EBITDA and cost categories · year-on-year", "3 / 7")
cats = ["EBITDA", "Materials", "Personnel", "Services", "Lease & rentals", "Other op. costs"]
colA = {"EBITDA":"84BD00","Materials":"F2A100","Personnel":"D9531E","Services":"3C7DA6","Lease & rentals":"8C5A2B","Other op. costs":"BFBFBF"}
txtA = {"EBITDA":INK,"Materials":INK,"Personnel":WHITE,"Services":WHITE,"Lease & rentals":WHITE,"Other op. costs":INK}
d23 = {"EBITDA":7.0,"Materials":46.0,"Personnel":23.1,"Services":19.9,"Lease & rentals":3.8,"Other op. costs":0.3}
d24 = {"EBITDA":3.8,"Materials":38.4,"Personnel":30.3,"Services":22.8,"Lease & rentals":4.2,"Other op. costs":0.5}
deltas = {"EBITDA":"-3.2","Materials":"-7.6","Personnel":"+7.2","Services":"+2.9","Lease & rentals":"+0.4","Other op. costs":"+0.2"}
TOPB, SCALE, BW = 165, 6.0, 130
def draw_bar(x, data):
    y = TOPB; bounds = [y]
    for c in cats:
        h = data[c]*SCALE
        rect(s, x, y, BW, h, fill=colA[c])
        if h >= 17:
            para(tb(s, x, y, BW, h, "m"), f"{data[c]:.1f}%", size=(13 if h >= 28 else 10), color=txtA[c], bold=True, align="c", first=True)
        y += h; bounds.append(y)
    return bounds
BX1, BX2 = 480, 720
b1 = draw_bar(BX1, d23)
b2 = draw_bar(BX2, d24)
for i in range(len(b1)):
    cn = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, PX(BX1+BW), PX(b1[i]), PX(BX2), PX(b2[i]))
    cn.line.color.rgb = RGBColor.from_string("D9DEE6"); cn.line.width = Pt(0.75); cn.shadow.inherit = False
para(tb(s, BX1-20, b1[-1]+8, BW+40, 24), "2023", size=15, color=BLUE, bold=True, align="c", first=True)
para(tb(s, BX2-20, b2[-1]+8, BW+40, 24), "2024", size=15, color=BLUE, bold=True, align="c", first=True)
para(tb(s, 350, TOPB-2, 120, 22), "% of value", size=11, color=MUTED, align="r", first=True)
para(tb(s, 350, TOPB+16, 120, 22), "of production", size=11, color=MUTED, align="r", first=True)
# recap table (manual, fully styled)
rect(s, 960, 190, 590, 300, fill=LIGHT, line=BORDER)
para(tb(s, 982, 200, 260, 26), "Category", size=12, color=GRAY, bold=True, first=True)
para(tb(s, 1250, 200, 96, 26), "2023", size=12, color=GRAY, bold=True, align="c", first=True)
para(tb(s, 1352, 200, 96, 26), "2024", size=12, color=GRAY, bold=True, align="c", first=True)
para(tb(s, 1454, 200, 90, 26), "Δ pp", size=12, color=GRAY, bold=True, align="c", first=True)
rect(s, 982, 230, 562, 1, fill=BORDER)
ry = 244
for c in cats:
    para(tb(s, 982, ry, 268, 28, "m"), [("■ ", False, colA[c], 12), (c, True, INK, 12)], first=True)
    para(tb(s, 1250, ry, 96, 28, "m"), f"{d23[c]:.1f}%", size=12, color=TXT, align="c", first=True)
    para(tb(s, 1352, ry, 96, 28, "m"), f"{d24[c]:.1f}%", size=12, color=TXT, align="c", first=True)
    dv = deltas[c]; dcol = "2E7D4F" if (dv.startswith("-") and c != "EBITDA") else "C8102E"
    para(tb(s, 1454, ry, 90, 28, "m"), dv, size=12, color=dcol, bold=True, align="c", first=True)
    ry += 36
# callout
rect(s, 960, 505, 590, 150, fill=BLUEBG, line=BLUELN)
tfc = tb(s, 982, 517, 546, 130)
para(tfc, [("Fixed-cost absorption is the story: ", True, BLUE, 14), ("personnel rises ", False, TXT, 14), ("+7pp", True, TXT, 14), (" of output (stable in €, heavier on lower volume) while materials fall ", False, TXT, 14), ("-8pp", True, TXT, 14), (" (variable) — net, EBITDA compresses from 7.0% to 3.8%.", False, TXT, 14)], first=True, after=6)
para(tfc, "Read with the project cycle: 2024 value of production is lower as the 2023 WIP build-up unwinds.", size=12, color=GRAY, italic=True)
para(tb(s, 50, 860, 1330, 30), "% of value of production (€52.0m 2023 / €40.7m 2024). 2023 base inflated by WIP build-up (+€19.3m); on a revenue basis EBITDA margin was 7.2% (2023) / 3.5% (2024). Source: Coface reclassified P&L.", size=9.5, color=MUTED, first=True)

# =====================================================================
# SLIDE 2b — CUSTOMER MAPPING (B2B end-markets)
# =====================================================================
s = slide()
header(s, "Who buys from Hennecke — B2B customer map",
       "Customers are industrial manufacturers that process polyurethane · by application", "4 / 7")

def seg_card(x, y, w, h, title, desc, examples, accent=BLUE):
    rect(s, x, y, w, h, fill=WHITE, line=BORDER)
    rect(s, x, y, w, 42, fill=accent)
    para(tb(s, x+16, y, w-22, 42, "m"), title, size=13, color=WHITE, bold=True, first=True)
    tf = tb(s, x+16, y+52, w-30, h-62)
    para(tf, desc, size=11.5, color=GRAY, italic=True, first=True, after=6)
    if examples:
        para(tf, [("e.g. ", True, INK, 11.5), (examples, False, INK, 11.5)], after=0)

COLX = [50, 430, 810, 1190]
R1, R2 = 168, 452
CW, CH = 360, 252
seg_card(COLX[0], R1, CW, CH, "1 · Building insulation — sandwich panels",
         "Continuous & discontinuous panel makers (roofing, façades, cold rooms).",
         "Kingspan, Isopan / Manni, Metecno, Brucha, Trimo, Tata Steel")
seg_card(COLX[1], R1, CW, CH, "2 · Refrigeration & cold chain",
         "Fridge / freezer & cold-room manufacturers (PU injection insulation).",
         "Whirlpool, Electrolux, BSH, Liebherr, Haier, Arçelik; Epta, Arneg")
seg_card(COLX[2], R1, CW, CH, "3 · Flexible foam (slabstock)",
         "Foam producers supplying bedding & upholstered furniture.",
         "Recticel, The Vita Group, Eurofoam, Carpenter, Orsa Foam")
seg_card(COLX[3], R1, CW, CH, "4 · Automotive",
         "Tier-1 makers of seating, interiors, NVH, RIM & composite parts.",
         "Adient, Lear, Forvia, Magna, Grupo Antolin")
seg_card(COLX[0], R2, CW, CH, "5 · Footwear",
         "PU sole producers & footwear makers (direct injection); strong Italian districts.",
         "Sole manufacturers & footwear brands' suppliers")
seg_card(COLX[1], R2, CW, CH, "6 · Elastomers & composites",
         "Technical PU parts: rollers, coatings, RIM and composite components.",
         "Industrial manufacturers of technical PU parts")
seg_card(COLX[2], R2, CW, CH, "7 · Pipe insulation & water heaters",
         "District-heating pre-insulated pipes; boiler / water-heater makers.",
         "Logstor (Kingspan), Brugg, Isoplus; Ariston, A.O. Smith")
seg_card(COLX[3], R2, CW, CH, "8 · PU system houses (partners)",
         "Chemical groups integrating process technology with raw-material systems.",
         "Covestro, BASF, Dow, Huntsman")
# bottom note band
rect(s, 50, 718, 1500, 96, fill=BLUEBG, line=BLUELN)
tfn = tb(s, 74, 728, 1452, 80, "m")
para(tfn, [("Read-across: ", True, BLUE, 14),
           ("capex-driven, cyclical buyers → the installed base drives recurring ", False, TXT, 14),
           ("aftermarket", True, TXT, 14),
           (" (spares, retrofit, service). Reach is Italy/Europe-centric, global via the Hennecke network (US, Asia).", False, TXT, 14)],
     first=True)
para(tb(s, 50, 866, 1300, 20),
     "Illustrative customer profiles by application — not a confirmed client list (customer breakdown not disclosed in the analysed documents). Company names are sector examples.",
     size=9.5, color=MUTED, first=True)

# =====================================================================
# SLIDE 2b2 — COMPETITOR MAP (positioning)
# =====================================================================
s = slide()
header(s, "Competitor map — PU machinery landscape",
       "Qualitative positioning of the main players in the reference sector", "5 / 7")
# --- plot frame ---
rect(s, 120, 158, 880, 530, fill="FBFCFD", line=BORDER)
PX0, PX1, PY0, PY1 = 150, 970, 190, 660  # data plot bounds (px)
def mx(xp): return PX0 + xp/100.0*(PX1-PX0)
def my(yp): return PY1 - yp/100.0*(PY1-PY0)
midx, midy = (PX0+PX1)/2, (PY0+PY1)/2
rect(s, PX0, midy, PX1-PX0, 1.4, fill="C9D2DD")   # x-axis
rect(s, midx, PY0, 1.4, PY1-PY0, fill="C9D2DD")   # y-axis
# axis descriptors
para(tb(s, PX0, 165, PX1-PX0, 20, "t"), "▲ Global scale & installed base", size=11.5, color=GRAY, bold=True, align="c", first=True)
para(tb(s, PX0, 665, PX1-PX0, 20, "t"), "▼ Regional / niche reach", size=11.5, color=GRAY, bold=True, align="c", first=True)
para(tb(s, 128, midy+6, 230, 18, "t"), "◄ Specialised PU pure-play", size=11, color=MUTED, align="l", first=True)
para(tb(s, PX1-240, midy+6, 230, 18, "t"), "Diversified machinery group ►", size=11, color=MUTED, align="r", first=True)

def bubble(name, xp, yp, color, r, tag=None, tagcolor=None):
    cx, cy = mx(xp), my(yp)
    sh = s.shapes.add_shape(MSO_SHAPE.OVAL, PX(cx-r), PX(cy-r), PX(2*r), PX(2*r))
    sh.shadow.inherit = False
    sh.fill.solid(); sh.fill.fore_color.rgb = RGBColor.from_string(color)
    sh.line.color.rgb = RGBColor.from_string(WHITE); sh.line.width = Pt(1.5)
    lab = tb(s, cx-90, cy+r+1, 180, 34, "t")
    para(lab, name, size=11.5, color=INK, bold=True, align="c", first=True, after=0, spacing=0.95)
    if tag:
        para(lab, tag, size=9.5, color=tagcolor or MUTED, align="c", italic=True, after=0)

# competitors (x%: specialised→diversified ; y%: regional→global)
bubble("Hennecke Group", 33, 86, BLUE, 27, "PU world leader (incl. OMS)", BLUE)
bubble("Hennecke-OMS", 21, 66, RED, 22, "TARGET", RED)
bubble("Cannon / Afros", 31, 50, "5B7A9E", 18)
bubble("KraussMaffei", 84, 82, "5B7A9E", 23, "RPM unit · ChemChina", MUTED)
bubble("Kurtz Ersa", 74, 54, "5B7A9E", 16)
bubble("Frimo", 60, 42, "5B7A9E", 16)
bubble("Saip", 20, 26, "5B7A9E", 14)

# --- right panel: key players ---
rect(s, 1020, 158, 530, 530, fill=WHITE, line=BORDER)
para(tb(s, 1044, 170, 480, 26), "Key players", size=16, color=BLUE, bold=True, first=True)
rect(s, 1044, 200, 482, 1, fill=BORDER)
tfp = tb(s, 1044, 212, 484, 460)
players = [
    [("Hennecke Group ", True, BLUE), ("(DE/IT) — world leader in PU process tech; includes the target Hennecke-OMS.", False, INK)],
    [("Cannon / Afros ", True, INK), ("(Peschiera Borromeo, IT) — PU metering, foam & composites; main Italian rival.", False, INK)],
    [("KraussMaffei ", True, INK), ("(Munich, DE) — Reaction Process Machinery unit; diversified, ChemChina-owned.", False, INK)],
    [("Frimo Group ", True, INK), ("(Lotte, DE) — PU tooling & systems, automotive-focused.", False, INK)],
    [("Kurtz Ersa ", True, INK), ("(Kreuzwertheim, DE) — EPS/PU foaming machinery.", False, INK)],
    [("Saip ", True, INK), ("(Offanengo, IT) — PU plants / refrigeration insulation.", False, INK)],
]
for i, pl in enumerate(players):
    para(tfp, [("•  ", False, "5B7A9E", 13)] + pl, size=13, first=(i == 0), after=13, spacing=1.0)

para(tb(s, 50, 700, 1000, 24, "t"),
     "Net read: the target sits among specialised PU pure-plays with global, group-backed reach — differentiated in sandwich-panel & flexible-foam lines.",
     size=12.5, color=BLUE, bold=True, first=True)
para(tb(s, 50, 866, 1300, 20),
     "Qualitative positioning — illustrative, not to scale; consolidated competitor revenues not disclosed in the analysed documents.",
     size=9.5, color=MUTED, first=True)

# =====================================================================
# SLIDE 2c — TAILWINDS vs HEADWINDS
# =====================================================================
s = slide()
header(s, "Tailwinds vs headwinds — what drives this business",
       "Two silos: structural supports vs structural pressures, specific to a PU capital-goods model", "6 / 7")
# left silo — TAILWINDS
rect(s, 50, 125, 740, 575, fill=WHITE, line=BORDER)
card_title(s, 50, 125, 740, "TAILWINDS  ·  structural supports", GREEN)
tft = tb(s, 74, 188, 692, 500)
bullets(tft, [
    [("Building energy-efficiency regulation", True), (" (EU EPBD, Green Deal, renovations) → demand for insulating sandwich panels & rigid foam.", False)],
    [("Cold-chain & refrigeration", True), (" expansion (food, pharma) → PU-injection lines; Hennecke already equips much of global fridge output.", False)],
    [("Blowing-agent transition", True), (" (HFC → HFO/pentane) forces a customer machine-fleet renewal cycle.", False)],
    [("e-mobility & lightweighting", True), (" → automotive PU seating, NVH and composite parts.", False)],
    [("Large global installed base", True), (" → recurring, high-margin aftermarket (spares, retrofit, service).", False)],
    [("Brückner Group backing", True), (" → R&D, cross-selling, global reach and financial strength.", False)],
    [("Historic OMS know-how", True), (" in continuous sandwich-panel & flexible-foam lines (differentiation).", False)],
], accent=GREEN, size=14.5, color=INK, after=10)
# right silo — HEADWINDS
rect(s, 810, 125, 740, 575, fill=WHITE, line=BORDER)
card_title(s, 810, 125, 740, "HEADWINDS  ·  structural pressures", RED)
tfh = tb(s, 834, 188, 692, 500)
bullets(tfh, [
    [("Capex-driven cyclicality", True), (" → lumpy, volatile order intake tied to customers' investment cycle.", False)],
    [("Long project cycles", True), (" → work-in-progress swings distort reported earnings (the 2024 margin effect).", False)],
    [("Thin, volatile margins", True), (" & execution risk on complex turnkey orders.", False)],
    [("Captive nature", True), (" → limited autonomy, intercompany dependence, carve-out complexity.", False)],
    [("Input-cost volatility", True), (" (steel, energy, components) and supply-chain tensions.", False)],
    [("FX & country risk", True), (" on export projects; qualified competition (Cannon/Afros, KraussMaffei) → price pressure.", False)],
    [("Order concentration", True), (" on single large projects; warranty / performance-bond exposure.", False)],
], accent=RED, size=14.5, color=INK, after=10)
# bottom synthesis band
rect(s, 50, 718, 1500, 96, fill=BLUEBG, line=BLUELN)
tfs = tb(s, 74, 728, 1452, 80, "m")
para(tfs, [("Net read:  ", True, BLUE, 14),
           ("structurally-supported demand on a cyclical, execution-heavy model. The thesis hinges on ", False, TXT, 14),
           ("smoothing the cycle", True, TXT, 14), (" (backlog & pricing discipline) and ", False, TXT, 14),
           ("scaling the recurring aftermarket", True, TXT, 14),
           (" to lift and stabilise margins.", False, TXT, 14)], first=True)
para(tb(s, 50, 866, 1300, 20),
     "Drivers synthesised from the dossier (sections 3, 5, 9, 11) and public industry sources.",
     size=9.5, color=MUTED, first=True)

# =====================================================================
# SLIDE 3 — THESIS
# =====================================================================
s = slide()
header(s, "Investment thesis — value, risks & recommendation",
       "Margin-recovery & aftermarket story on a structurally-supported PU platform", "7 / 7")
# Column A — value creation
rect(s, 50, 125, 480, 560, fill=WHITE, line=BORDER)
card_title(s, 50, 125, 480, "Value-creation levers", GREEN)
tfa = tb(s, 74, 186, 432, 330)
bullets(tfa, [
    [("Margin recovery 3.5% → 6-8%", True), (" via fixed-cost efficiency.", False)],
    [("Grow the recurring ", False), ("aftermarket", True), (" (service, spares, retrofit).", False)],
    [("Working-capital optimisation (already underway: DIO 119→66).", False)],
    [("Backlog normalisation & pricing discipline on projects.", False)],
    [("Optional ", False), ("buy-and-build", True), (" in PU / adjacent capital-goods machinery.", False)],
], accent=GREEN, size=15, after=11)
rect(s, 74, 535, 432, 128, fill=GREENBG, line=GREENLN)
tfa2 = tb(s, 94, 547, 392, 110)
para(tfa2, "Structural demand drivers", size=14, color=GREEN, bold=True, first=True, after=5)
para(tfa2, "Building energy efficiency (insulating panels) · cold chain · e-mobility · blowing-agent transition (HFO/pentane).", size=14, color=TXT)
# Column B — risks
rect(s, 560, 125, 480, 560, fill=WHITE, line=BORDER)
card_title(s, 560, 125, 480, "Key risks & cautions", RED)
tfb = tb(s, 584, 186, 432, 330)
bullets(tfb, [
    [("Thin, volatile margins", True), ("; 2024 operating loss (-€1.0M EBIT).", False)],
    [("Captive nature", True), (" → carve-out complexity (brand, IT, shared svcs).", False)],
    [("Heavy ", False), ("intercompany", True), (" ties & group treasury — limited autonomy.", False)],
    [("Pledged shares", True), (" to be released; thin equity (€6.0M).", False)],
    [("Intangibles €12.0M", True), (" — impairment risk; cyclicality & execution risk.", False)],
], accent=RED, size=15, after=11)
rect(s, 584, 535, 432, 128, fill=PINKBG, line=PINKLN)
tfb2 = tb(s, 604, 547, 392, 110)
para(tfb2, "Mitigants", size=14, color=RED, bold=True, first=True, after=5)
para(tfb2, "Light balance sheet & high liquidity; global installed base; world-class group R&D and sales network behind the asset.", size=14, color=TXT)
# Column C — valuation & next steps
rect(s, 1070, 125, 480, 560, fill=WHITE, line=BORDER)
card_title(s, 1070, 125, 480, "Valuation & next steps", BLUE)
rect(s, 1094, 192, 432, 92, fill=BLUEBG, line=BLUELN)
tfc0 = tb(s, 1110, 200, 400, 78, "m")
para(tfc0, "Indicative EV/EBITDA", size=14, color=GRAY, first=True, after=2)
para(tfc0, [("6–9x", True, BLUE, 28), ("   (5–7x small-cap)", False, MUTED, 15)], after=0)
tfc = tb(s, 1094, 296, 432, 90)
para(tfc, [("Value on ", False, TXT), ("normalised EBITDA", True, TXT),
           (", not the depressed 2024 figure; consider asset value.", False, TXT)], size=14, first=True)
para(tb(s, 1094, 380, 432, 26), "Proceed to focused Due Diligence:", size=15, color=BLUE, bold=True, first=True)
tfc2 = tb(s, 1094, 408, 432, 160)
bullets(tfc2, [
    "Transfer pricing & intercompany flows",
    "Backlog & project margins",
    "Carve-out feasibility (brand / IT / people)",
    "EBITDA normalisation & true net debt",
    "Share-pledge release & intangibles testing",
], accent=BLUE, size=14, color=INK, after=7)
rect(s, 1094, 585, 432, 78, fill=BLUE)
tfc3 = tb(s, 1094, 593, 432, 64, "m")
para(tfc3, "RECOMMENDATION", size=13, color="CDD8E6", align="c", first=True, after=3)
para(tfc3, "Approve progression to Phase-2 DD", size=17, color=WHITE, bold=True, align="c")
# bottom line band
rect(s, 50, 705, 1500, 120, fill=LIGHT, line=BORDER)
tfbl = tb(s, 74, 716, 1452, 105)
para(tfbl, "Bottom line", size=16, color=BLUE, bold=True, first=True, after=5)
para(tfbl, [("A historic-brand PU engineering platform with a global installed base and structural tailwinds, currently held back by cyclical, thin margins. The opportunity is a ", False, TXT),
            ("margin-normalisation & aftermarket", True, TXT),
            (" play — execution and carve-out feasibility are the swing factors.", False, TXT)], size=15)
para(tb(s, 50, 866, 1000, 20), "For Board discussion — not an investment recommendation; subject to specialist Due Diligence.",
     size=9.5, color=MUTED, first=True)

# =====================================================================
# SLIDE 4 — CLOSING / Q&A
# =====================================================================
s = slide()
rect(s, 0, 0, 1600, 900, fill=BLUE)
rect(s, 0, 0, 1600, 10, fill=RED)
para(tb(s, 120, 225, 600, 28), "DISCUSSION", size=15, color=LBLUE, first=True)
rect(s, 120, 272, 90, 5, fill=RED)
para(tb(s, 116, 300, 1200, 110), "Questions & discussion", size=62, color=WHITE, bold=True, first=True)
# recommendation banner
rect(s, 120, 430, 900, 92, fill=RED)
tfrb = tb(s, 148, 442, 850, 70, "m")
para(tfrb, "RECOMMENDATION", size=13, color="FFD7DC", first=True, after=3)
para(tfrb, "Approve progression to Phase-2 Due Diligence", size=24, color=WHITE, bold=True)
# key questions
para(tb(s, 120, 560, 800, 30), "Key questions to align on", size=20, color=GOLD, bold=True, first=True)
tfq = tb(s, 120, 600, 940, 220)
bullets(tfq, [
    "Preferred structure: carve-out of the Italian entity vs whole-Group platform?",
    "Risk appetite on the captive / intercompany & carve-out complexity?",
    "Valuation basis — normalised EBITDA assumptions and entry multiple?",
    "Budget & timeline for Phase-2 DD (financial, commercial, legal, tax)?",
], accent=RED, size=18, color="E3E9F1", after=12)
# next step card
rect(s, 1090, 430, 390, 345, fill=PANEL, line="3A5578")
para(tb(s, 1115, 445, 350, 30), "Next step", size=17, color=WHITE, bold=True, first=True)
rect(s, 1115, 482, 340, 1, fill="3A5578")
tfns = tb(s, 1115, 496, 350, 140)
para(tfns, "Mandate the deal team to launch focused Due Diligence and return with a normalised valuation and indicative offer range.",
     size=15, color="CDD8E6", first=True)
rect(s, 1115, 636, 340, 2, fill="3A5578")
tfns2 = tb(s, 1115, 656, 350, 100)
para(tfns2, "Project lead", size=14, color=LBLUE, first=True, after=4)
para(tfns2, "Corporate Development / M&A", size=17, color=WHITE, bold=True, after=6)
para(tfns2, "Materials available on request", size=13, color="7E93AC")
rect(s, 120, 822, 1360, 1, fill="2D4769")
para(tb(s, 120, 834, 1000, 22), "CONFIDENTIAL — For Board discussion only. Not an investment recommendation; subject to specialist Due Diligence.",
     size=11, color="7E93AC", first=True)
para(tb(s, 1080, 834, 400, 22, "t"), "HENNECKE-OMS S.p.A. · June 2026", size=11, color="7E93AC", align="r", first=True)

prs.save(OUT)
print("Saved editable deck:", OUT, "| slides:", len(prs.slides._sldIdLst))
