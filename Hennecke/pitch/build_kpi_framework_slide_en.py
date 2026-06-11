#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Standalone editable KPI-framework slide (native PowerPoint)."""
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

OUT = "/home/user/Storage/Hennecke/pitch/kpi_framework_slide_EN.pptx"
BLUE="1F3A5F"; RED="C8102E"; GREEN="2E7D4F"; GOLD="B8860B"; WHITE="FFFFFF"
GRAY="666666"; MUTED="999999"; INK="222222"; BORDER="DFE4EA"
GBG="F7FBF8"; GLN="CFE6D8"; BBG="F5F8FC"; BLN="CDD8E6"; ABG="FBF6EA"; ALN="EAD9B0"; PANEL="EEF2F7"

EMU=7620
def PX(v): return Emu(int(round(v*EMU)))
prs=Presentation(); prs.slide_width=Emu(int(13.333*914400)); prs.slide_height=Emu(int(7.5*914400))
s=prs.slides.add_slide(prs.slide_layouts[6])
ALIGN={"l":PP_ALIGN.LEFT,"c":PP_ALIGN.CENTER,"r":PP_ALIGN.RIGHT}
ANCH={"t":MSO_ANCHOR.TOP,"m":MSO_ANCHOR.MIDDLE,"b":MSO_ANCHOR.BOTTOM}

def rect(x,y,w,h,fill=None,line=None,rot=0):
    sh=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,PX(x),PX(y),PX(w),PX(h)); sh.shadow.inherit=False
    if fill is None: sh.fill.background()
    else: sh.fill.solid(); sh.fill.fore_color.rgb=RGBColor.from_string(fill)
    if line is None: sh.line.fill.background()
    else: sh.line.color.rgb=RGBColor.from_string(line); sh.line.width=Pt(1)
    if rot: sh.rotation=rot
    return sh

def tb(x,y,w,h,anchor="t",rot=0):
    b=s.shapes.add_textbox(PX(x),PX(y),PX(w),PX(h)); tf=b.text_frame; tf.word_wrap=True
    tf.vertical_anchor=ANCH[anchor]
    tf.margin_left=Pt(3); tf.margin_right=Pt(3); tf.margin_top=Pt(1); tf.margin_bottom=Pt(1)
    if rot: b.rotation=rot
    return tf

def para(tf,runs,size=13,color=INK,bold=False,align="l",first=False,after=5):
    p=tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment=ALIGN[align]; p.space_after=Pt(after); p.space_before=Pt(0)
    if isinstance(runs,str): runs=[(runs,bold,color,size)]
    for rn in runs:
        r=p.add_run(); r.text=rn[0]; f=r.font; f.name="Arial"
        f.size=Pt(rn[3] if len(rn)>3 and rn[3] else size)
        f.bold=rn[1] if len(rn)>1 and rn[1] is not None else bold
        f.color.rgb=RGBColor.from_string(rn[2] if len(rn)>2 and rn[2] else color)
    return p

# header
rect(0,0,1600,8,fill=RED)
para(tb(50,28,1200,40),"KPI framework — growth, profitability & cash",size=26,color=BLUE,bold=True,first=True)
para(tb(50,76,1200,26),"Commercial & Operations levers underpinning the monitoring of the target's business",size=13,color=GRAY,first=True)
para(tb(1150,26,400,22),"CONFIDENTIAL · For Board discussion",size=9.5,color=MUTED,align="r",first=True)

# column headers
for x,txt,col in [(210,"GROWTH",GREEN),(670,"PROFITABILITY",BLUE),(1130,"CASH",GOLD)]:
    rect(x,118,430,46,fill=col)
    para(tb(x,118,430,46,"m"),txt,size=16,color=WHITE,bold=True,align="c",first=True)

# row labels (rotated)
para(tb(20,250,150,40,"m",rot=270),"COMMERCIAL",size=15,color=RED,bold=True,align="c",first=True)
para(tb(20,535,150,40,"m",rot=270),"OPERATIONS",size=15,color=RED,bold=True,align="c",first=True)

cells = {
 (0,0):["Order intake & Book-to-bill","Backlog (months of coverage)","Win rate / hit rate on bids","Aftermarket share %","Application & geography mix"],
 (0,1):["Bid margin per project","% orders above margin threshold","Pricing waterfall / discount leakage","Price / Volume / Mix variance","Aftermarket vs new-machine margin"],
 (0,2):["Down-payment % at order","Billing-milestone (SAL) adherence","DSO (days sales outstanding)","Advances ÷ WIP","Penalties / performance bonds"],
 (1,0):["Capacity utilisation (shop floor & eng.)","On-Time Delivery (OTD)","Order-to-delivery lead time","Throughput (machines completed)","NPD / time-to-market"],
 (1,1):["EAC vs budget / margin-at-completion","Productivity (std vs actual hrs) / OEE","Scrap · give-away · material yield","Material cost ratio & PPV","Warranty cost · revenue/FTE"],
 (1,2):["DIO inventory (days)","Unbilled WIP & over/under-billing","DPO (days payable)","Cash Conversion Cycle","Capex vs plan · FCF conversion"],
}
colx=[210,670,1130]; rowy=[178,478]
fills=[(GBG,GLN),(BBG,BLN),(ABG,ALN)]
for r in (0,1):
    for c in (0,1,2):
        x=colx[c]; y=rowy[r]; fb,ln=fills[c]
        rect(x,y,430,288,fill=fb,line=ln)
        tf=tb(x+16,y+18,402,260)
        for i,k in enumerate(cells[(r,c)]):
            para(tf,[("•  ",False,(GREEN if c==0 else BLUE if c==1 else GOLD)),(k,False,INK)],size=13,first=(i==0),after=14)

# bottom band
rect(50,784,1510,66,fill=PANEL,line=BLN)
tfb=tb(70,790,1480,58)
para(tfb,[("Captive lens: ",True,BLUE,13),("intercompany % (purchases/sales) & stand-alone margin — to measure true profitability and prepare a possible carve-out.",False,INK,13)],first=True,after=3)
para(tfb,[("CFO cockpit (10): ",True,BLUE,12.5),("order intake · book-to-bill · backlog · aftermarket % · normalised EBITDA · margin-at-completion · OTD · CCC · advances÷WIP · FCF conversion.",False,"333333",12.5)],after=0)

prs.save(OUT); print("Saved:",OUT)
