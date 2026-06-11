#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Extract ONLY the competitor-map slide into a standalone editable PPTX."""
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

OUT = "/home/user/Storage/Hennecke/pitch/competitor_map_slide.pptx"
BLUE="1F3A5F"; RED="C8102E"; WHITE="FFFFFF"; GRAY="666666"; MUTED="999999"
BORDER="DFE4EA"; INK="222222"; TXT="333333"; BLUEBG="EEF2F7"; BLUELN="CDD8E6"

EMU = 7620
def PX(v): return Emu(int(round(v*EMU)))
prs = Presentation()
prs.slide_width = Emu(int(13.333*914400)); prs.slide_height = Emu(int(7.5*914400))
s = prs.slides.add_slide(prs.slide_layouts[6])
ALIGN = {"l":PP_ALIGN.LEFT, "c":PP_ALIGN.CENTER, "r":PP_ALIGN.RIGHT}
ANCH = {"t":MSO_ANCHOR.TOP, "m":MSO_ANCHOR.MIDDLE, "b":MSO_ANCHOR.BOTTOM}

def rect(x, y, w, h, fill=None, line=None):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, PX(x), PX(y), PX(w), PX(h))
    sh.shadow.inherit = False
    if fill is None: sh.fill.background()
    else: sh.fill.solid(); sh.fill.fore_color.rgb = RGBColor.from_string(fill)
    if line is None: sh.line.fill.background()
    else: sh.line.color.rgb = RGBColor.from_string(line); sh.line.width = Pt(1)
    return sh

def tb(x, y, w, h, anchor="t"):
    box = s.shapes.add_textbox(PX(x), PX(y), PX(w), PX(h)); tf = box.text_frame
    tf.word_wrap = True; tf.vertical_anchor = ANCH[anchor]
    tf.margin_left = Pt(2); tf.margin_right = Pt(2); tf.margin_top = Pt(1); tf.margin_bottom = Pt(1)
    return tf

def para(tf, runs, size=14, color=INK, bold=False, align="l", first=False, after=4, italic=False, spacing=1.0):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = ALIGN[align]; p.space_after = Pt(after); p.space_before = Pt(0)
    try: p.line_spacing = spacing
    except Exception: pass
    if isinstance(runs, str): runs = [(runs, bold, color, size)]
    for rn in runs:
        b = rn[1] if len(rn) > 1 and rn[1] is not None else bold
        c = rn[2] if len(rn) > 2 and rn[2] is not None else color
        sz = rn[3] if len(rn) > 3 and rn[3] is not None else size
        r = p.add_run(); r.text = rn[0]; f = r.font
        f.name = "Arial"; f.size = Pt(sz); f.bold = b; f.italic = italic
        f.color.rgb = RGBColor.from_string(c)
    return p

# header
rect(0, 0, 1600, 8, fill=RED)
para(tb(50, 30, 1100, 50), "Competitor map — PU machinery landscape", size=26, color=BLUE, bold=True, first=True)
para(tb(50, 78, 1200, 30), "Qualitative positioning of the main players in the reference sector", size=13, color=GRAY, first=True)
para(tb(1150, 26, 400, 22), "CONFIDENTIAL · For Board discussion", size=9.5, color=MUTED, align="r", first=True)
rect(50, 858, 1500, 1, fill=BORDER)
para(tb(1400, 866, 150, 20), "4 / 6", size=9.5, color=BLUE, bold=True, align="r", first=True)

# plot
rect(120, 158, 880, 530, fill="FBFCFD", line=BORDER)
PX0, PX1, PY0, PY1 = 150, 970, 190, 660
def mx(xp): return PX0 + xp/100.0*(PX1-PX0)
def my(yp): return PY1 - yp/100.0*(PY1-PY0)
midx, midy = (PX0+PX1)/2, (PY0+PY1)/2
rect(PX0, midy, PX1-PX0, 1.4, fill="C9D2DD"); rect(midx, PY0, 1.4, PY1-PY0, fill="C9D2DD")
para(tb(PX0, 165, PX1-PX0, 20), "▲ Global scale & installed base", size=11.5, color=GRAY, bold=True, align="c", first=True)
para(tb(PX0, 665, PX1-PX0, 20), "▼ Regional / niche reach", size=11.5, color=GRAY, bold=True, align="c", first=True)
para(tb(128, midy+6, 230, 18), "◄ Specialised PU pure-play", size=11, color=MUTED, align="l", first=True)
para(tb(PX1-240, midy+6, 230, 18), "Diversified machinery group ►", size=11, color=MUTED, align="r", first=True)

def bubble(name, xp, yp, color, r, tag=None, tagcolor=None):
    cx, cy = mx(xp), my(yp)
    sh = s.shapes.add_shape(MSO_SHAPE.OVAL, PX(cx-r), PX(cy-r), PX(2*r), PX(2*r))
    sh.shadow.inherit = False; sh.fill.solid(); sh.fill.fore_color.rgb = RGBColor.from_string(color)
    sh.line.color.rgb = RGBColor.from_string(WHITE); sh.line.width = Pt(1.5)
    lab = tb(cx-90, cy+r+1, 180, 34)
    para(lab, name, size=11.5, color=INK, bold=True, align="c", first=True, after=0, spacing=0.95)
    if tag: para(lab, tag, size=9.5, color=tagcolor or MUTED, align="c", italic=True, after=0)

bubble("Hennecke Group", 33, 86, BLUE, 27, "PU world leader (incl. OMS)", BLUE)
bubble("Hennecke-OMS", 21, 66, RED, 22, "TARGET", RED)
bubble("Cannon / Afros", 31, 50, "5B7A9E", 18)
bubble("KraussMaffei", 84, 82, "5B7A9E", 23, "RPM unit · ChemChina", MUTED)
bubble("Kurtz Ersa", 74, 54, "5B7A9E", 16)
bubble("Frimo", 60, 42, "5B7A9E", 16)
bubble("Saip", 20, 26, "5B7A9E", 14)

# right panel
rect(1020, 158, 530, 530, fill=WHITE, line=BORDER)
para(tb(1044, 170, 480, 26), "Key players", size=16, color=BLUE, bold=True, first=True)
rect(1044, 200, 482, 1, fill=BORDER)
tfp = tb(1044, 212, 484, 460)
players = [
    [("Hennecke Group ", True, BLUE), ("(DE/IT) — world leader in PU process tech; includes the target Hennecke-OMS.", False, INK)],
    [("Cannon / Afros ", True, INK), ("(Peschiera Borromeo, IT) — PU metering, foam & composites; main Italian rival.", False, INK)],
    [("KraussMaffei ", True, INK), ("(Munich, DE) — Reaction Process Machinery unit; diversified, ChemChina-owned.", False, INK)],
    [("Frimo Group ", True, INK), ("(Lotte, DE) — PU tooling & systems, automotive-focused.", False, INK)],
    [("Kurtz Ersa ", True, INK), ("(Kreuzwertheim, DE) — EPS/PU foaming machinery.", False, INK)],
    [("Saip ", True, INK), ("(Offanengo, IT) — PU plants / refrigeration insulation.", False, INK)],
]
for i, pl in enumerate(players):
    para(tfp, [("•  ", False, "5B7A9E", 13)] + pl, size=13, first=(i == 0), after=13)
para(tb(50, 700, 1000, 24), "Net read: the target sits among specialised PU pure-plays with global, group-backed reach — differentiated in sandwich-panel & flexible-foam lines.",
     size=12.5, color=BLUE, bold=True, first=True)
para(tb(50, 720, 1300, 20), "Qualitative positioning — illustrative, not to scale; consolidated competitor revenues not disclosed in the analysed documents.",
     size=9.5, color=MUTED, first=True)

prs.save(OUT)
print("Saved:", OUT)
